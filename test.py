from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

# PPT 생성
prs = Presentation()

# 슬라이드 추가 (빈 레이아웃)
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)

# 단계별 위치 정의 (좌표는 PPT 내 비율)
positions = {
    'start': (4, 0.5),
    'check_member': (4, 1.5),
    'register_form': (1, 2.5),
    'login_form': (7, 2.5),
    'user_info': (1, 3.5),
    'validate': (1, 4.5),
    'id_pw_check': (7, 3.5),
    'register_success': (1, 5.5),
    'login_success': (7, 4.5),
    'end': (4, 6.5),
}

# 도형 추가 함수 정의
def add_shape(slide, text, shape, left, top, width=Inches(2), height=Inches(0.7)):
    shape_obj = slide.shapes.add_shape(shape, Inches(left), Inches(top), width, height)
    shape_obj.text = text
    text_frame = shape_obj.text_frame
    p = text_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.name = 'NanumGothic'
    return shape_obj

# 단계 추가
shapes = {
    'start': add_shape(slide, '회원가입/로그인', MSO_SHAPE.OVAL, *positions['start']),
    'check_member': add_shape(slide, '회원 여부 판단', MSO_SHAPE.DIAMOND, *positions['check_member']),
    'register_form': add_shape(slide, '회원가입 폼 입력', MSO_SHAPE.RECTANGLE, *positions['register_form']),
    'login_form': add_shape(slide, '로그인 폼 입력', MSO_SHAPE.PARALLELOGRAM, *positions['login_form']),
    'user_info': add_shape(slide, '회원 정보 입력', MSO_SHAPE.PARALLELOGRAM, *positions['user_info']),
    'validate': add_shape(slide, '유효성 체크', MSO_SHAPE.DIAMOND, *positions['validate']),
    'id_pw_check': add_shape(slide, 'ID/PW 체크', MSO_SHAPE.DIAMOND, *positions['id_pw_check']),
    'register_success': add_shape(slide, '회원가입 성공', MSO_SHAPE.RECTANGLE, *positions['register_success']),
    'login_success': add_shape(slide, '로그인 처리', MSO_SHAPE.RECTANGLE, *positions['login_success']),
    'end': add_shape(slide, '메인화면 이동', MSO_SHAPE.OVAL, *positions['end']),
}

# 연결선 추가 함수
def connect_shapes(slide, shape1, shape2):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        shape1.left + shape1.width / 2,
        shape1.top + shape1.height,
        shape2.left + shape2.width / 2,
        shape2.top
    )
    connector.line.color.rgb = RGBColor(0, 0, 0)

# 흐름 연결
connect_shapes(slide, shapes['start'], shapes['check_member'])
connect_shapes(slide, shapes['check_member'], shapes['register_form'])
connect_shapes(slide, shapes['check_member'], shapes['login_form'])
connect_shapes(slide, shapes['register_form'], shapes['user_info'])
connect_shapes(slide, shapes['user_info'], shapes['validate'])
connect_shapes(slide, shapes['validate'], shapes['register_success'])
connect_shapes(slide, shapes['validate'], shapes['user_info'])
connect_shapes(slide, shapes['login_form'], shapes['id_pw_check'])
connect_shapes(slide, shapes['id_pw_check'], shapes['login_success'])
connect_shapes(slide, shapes['id_pw_check'], shapes['login_form'])
connect_shapes(slide, shapes['register_success'], shapes['end'])
connect_shapes(slide, shapes['login_success'], shapes['end'])

# 파일 저장
pptx_path = './app_flowchart_ppt_korean_fixed.pptx'
prs.save(pptx_path)

pptx_path